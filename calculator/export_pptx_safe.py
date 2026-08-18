from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

root = Path(__file__).resolve().parents[1]
output = root / "exports" / "pptx" / "Yv-Me_Investor_Donor_Grant_Presentation.pptx"

prs = Presentation()

def add_slide(title, lines):

    slide = prs.slides.add_slide(
        prs.slide_layouts[5]
    )

    slide.shapes.title.text = title

    box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.7),
        Inches(11),
        Inches(4.5)
    )

    tf = box.text_frame
    tf.clear()

    for index, line in enumerate(lines):

        paragraph = (
            tf.paragraphs[0]
            if index == 0
            else tf.add_paragraph()
        )

        paragraph.text = line
        paragraph.font.size = Pt(22)
        paragraph.space_after = Pt(14)

add_slide(
    "Yv-Me ROI Calculator",
    [
        "EasePal Care — Yobe State, Nigeria",
        "Programme Economics & Partnership ROI",
        "1 CHEW : 10 Beneficiaries"
    ]
)

add_slide(
    "Programme Snapshot",
    [
        "CHEWs: 100",
        "Beneficiaries: 1,000",
        "Programme Duration: 12 months",
        "Programme Budget: NGN 426,906,031"
    ]
)

add_slide(
    "Care Delivery Model",
    [
        "1 CHEW : 10 Beneficiaries",
        "Community-based care delivery",
        "Scalable deployment model"
    ]
)

add_slide(
    "Programme Economics",
    [
        "Monthly Programme Cost: NGN 35,575,503",
        "Annual Programme Cost: NGN 426,906,031",
        "Cost per CHEW: NGN 4,269,060",
        "Cost per Beneficiary: NGN 426,906"
    ]
)

add_slide(
    "Investor / Partner Opportunity",
    [
        "Adjustable programme assumptions",
        "Revenue scenario modelling",
        "ROI and break-even analysis",
        "Scaling economics"
    ]
)

add_slide(
    "Donor / Grant / iNGO Opportunity",
    [
        "Transparent programme budgeting",
        "Cost-per-beneficiary visibility",
        "Scalable care delivery",
        "Sustainability analysis"
    ]
)

add_slide(
    "Next Steps",
    [
        "Validate assumptions",
        "Select financing scenario",
        "Run scale scenarios",
        "Generate final investment or grant case"
    ]
)

prs.save(output)

print("PPTX CREATED SUCCESSFULLY")
print(output)
