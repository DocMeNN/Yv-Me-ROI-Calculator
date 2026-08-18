from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from calculator.analytics import calculate

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "exports" / "pptx" / "Yv-Me_Investor_Donor_Grant_Presentation.pptx"

CHEWS = 100
RATIO = 10
MONTHS = 12
BUDGET = 426906031
SUBSCRIPTION = 15000
INVESTMENT = 426906031
REVENUE_SHARE = 10

r = calculate(
    chews=CHEWS,
    beneficiaries_per_chew=RATIO,
    programme_months=MONTHS,
    total_budget=BUDGET,
    subscription_per_beneficiary=SUBSCRIPTION,
    investment_amount=INVESTMENT,
    revenue_share=REVENUE_SHARE,
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def slide(title, bullets):

    layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)

    title_box = s.shapes.title
    title_box.text = title

    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True

    box = s.shapes.add_textbox(
        Inches(0.8),
        Inches(1.5),
        Inches(11.7),
        Inches(5.2)
    )

    tf = box.text_frame
    tf.clear()

    for i, text in enumerate(bullets):

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        p.text = text
        p.font.size = Pt(20)
        p.space_after = Pt(12)

    return s

# ------------------------------------------------------------
# SLIDES
# ------------------------------------------------------------

slide(
    "Yv-Me ROI Calculator",
    [
        "EasePal Care — Yobe State, Nigeria",
        "Programme Economics & Partnership Investment Model",
        "Care Delivery Model: 1 CHEW : 10 Beneficiaries",
    ]
)

slide(
    "Programme Snapshot",
    [
        f"CHEWs: {CHEWS:,}",
        f"Beneficiaries: {r['beneficiaries']:,}",
        f"Programme Duration: {MONTHS} months",
        f"Programme Budget: NGN {BUDGET:,.0f}",
    ]
)

slide(
    "Care Delivery Model",
    [
        f"1 CHEW : {RATIO} Beneficiaries",
        f"{CHEWS:,} CHEWs support {r['beneficiaries']:,} beneficiaries",
        "Community-based care delivery model",
        "Scalable deployment architecture",
    ]
)

slide(
    "Programme Economics",
    [
        f"Monthly Programme Cost: NGN {r['monthly_operating_cost']:,.0f}",
        f"Annual Programme Cost: NGN {r['annual_operating_cost']:,.0f}",
        f"Cost per CHEW: NGN {r['cost_per_chew']:,.0f}",
        f"Cost per Beneficiary: NGN {r['cost_per_beneficiary']:,.0f}",
    ]
)

slide(
    "Revenue Opportunity",
    [
        f"Subscription Assumption: NGN {SUBSCRIPTION:,.0f} / beneficiary / month",
        f"Monthly Gross Revenue: NGN {r['monthly_gross_revenue']:,.0f}",
        f"Annual Gross Revenue: NGN {r['annual_gross_revenue']:,.0f}",
        "Revenue assumptions are adjustable in the calculator.",
    ]
)

slide(
    "ROI & Break-even",
    [
        f"Annual Contribution: NGN {r['annual_contribution']:,.0f}",
        f"ROI: {r['roi']:.2%}",
        (
            "Break-even: Not reached"
            if r["break_even_months"] == float("inf")
            else f"Break-even: {r['break_even_months']:.1f} months"
        ),
    ]
)

slide(
    "Scaling Opportunity",
    [
        "10 CHEWs → 100 beneficiaries",
        "100 CHEWs → 1,000 beneficiaries",
        "500 CHEWs → 5,000 beneficiaries",
        "1,000 CHEWs → 10,000 beneficiaries",
        "5,000 CHEWs → 50,000 beneficiaries",
    ]
)

slide(
    "Investor / Private Partnership Opportunity",
    [
        f"Illustrative Investment: NGN {INVESTMENT:,.0f}",
        f"Illustrative Revenue Share: {REVENUE_SHARE:.1f}%",
        "Scalable community healthcare delivery platform",
        "Adjustable commercial and investment assumptions",
    ]
)

slide(
    "Donor / Grant / iNGO Opportunity",
    [
        "Transparent programme budgeting",
        "Cost-per-beneficiary visibility",
        "Scalable deployment model",
        "Adjustable programme assumptions",
        "Evidence-ready financial analytics",
    ]
)

slide(
    "Decision Dashboard",
    [
        "Programme cost",
        "Cost per beneficiary",
        "Revenue opportunity",
        "ROI",
        "Break-even",
        "Scaling economics",
        "Investor / donor scenarios",
    ]
)

slide(
    "Next Steps",
    [
        "Validate programme assumptions",
        "Select financing scenario",
        "Confirm revenue / funding model",
        "Run scale scenarios",
        "Generate final investment or grant case",
    ]
)

prs.save(OUTPUT)

print(f"PPTX CREATED: {OUTPUT}")
